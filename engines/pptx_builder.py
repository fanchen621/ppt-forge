"""
PPT Forge v2 — PPTX Builder Engine (ENHANCED)
Assembles professional PPTX files from config + generated art.
Supports 20+ slide types including Chinese educational content.
Includes decorative shapes, mind maps, task separators, and more.
"""

import io
import random
from typing import List, Dict, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from PIL import Image, ImageFilter, ImageDraw, ImageEnhance
import numpy as np

from .art_engine import generate_art, get_palette, palette_color, adjust_brightness, with_alpha
from .layout_engine import SmartLayout, TYPE_SCALE
from .effects_pipeline import apply_effects
from .illustration_engine import auto_illustrate_slide, generate_illustration
from .svg_engine import svg_to_png as _svg_to_png
from .shapes_engine import (
    add_thought_bubble, add_framed_content, add_bracket_group,
    add_mind_map, add_mind_map_node, add_mind_map_connector,
    add_task_separator, add_vertical_text, add_number_badge,
    add_right_arrow_shape, add_card, add_step_box,
    add_styled_table
)


# ─── Theme Definitions (ENHANCED with Chinese Education Themes) ────────────

THEMES: Dict[str, Dict] = {
    # ─── Original Themes ──────────────────────────────────────────────
    "minimal": {
        "bg_color": (255, 255, 255),
        "title_color": (25, 25, 28),
        "body_color": (65, 65, 75),
        "accent_color": (0, 90, 200),
        "accent2_color": (80, 140, 240),
        "font_title": "Arial",
        "font_body": "Arial",
        "title_size": 44,
        "body_size": 16,
        "default_art": "gradient_mesh",
        "palette": "mono",
        "overlay_opacity": 0.05,
    },
    "dark": {
        "bg_color": (14, 14, 20),
        "title_color": (245, 245, 250),
        "body_color": (175, 175, 195),
        "accent_color": (99, 179, 237),
        "accent2_color": (56, 161, 240),
        "font_title": "Arial",
        "font_body": "Arial",
        "title_size": 44,
        "body_size": 16,
        "default_art": "flow_field",
        "palette": "midnight",
        "overlay_opacity": 0.12,
    },
    "gradient": {
        "bg_color": (25, 0, 45),
        "title_color": (255, 255, 255),
        "body_color": (220, 200, 240),
        "accent_color": (255, 100, 150),
        "accent2_color": (200, 80, 255),
        "font_title": "Arial",
        "font_body": "Arial",
        "title_size": 44,
        "body_size": 16,
        "default_art": "gradient_mesh",
        "palette": "cosmic",
        "overlay_opacity": 0.10,
    },
    "corporate": {
        "bg_color": (245, 247, 250),
        "title_color": (15, 45, 90),
        "body_color": (55, 75, 105),
        "accent_color": (0, 85, 175),
        "accent2_color": (0, 130, 220),
        "font_title": "Arial",
        "font_body": "Arial",
        "title_size": 40,
        "body_size": 16,
        "default_art": "wave",
        "palette": "cool",
        "overlay_opacity": 0.06,
    },
    "creative": {
        "bg_color": (255, 248, 238),
        "title_color": (35, 15, 5),
        "body_color": (85, 55, 35),
        "accent_color": (255, 90, 70),
        "accent2_color": (255, 140, 50),
        "font_title": "Arial",
        "font_body": "Arial",
        "title_size": 44,
        "body_size": 16,
        "default_art": "geometric",
        "palette": "warm",
        "overlay_opacity": 0.08,
    },
    "tech": {
        "bg_color": (6, 10, 18),
        "title_color": (0, 255, 200),
        "body_color": (140, 195, 215),
        "accent_color": (0, 200, 255),
        "accent2_color": (120, 0, 255),
        "font_title": "Arial",
        "font_body": "Arial",
        "title_size": 44,
        "body_size": 16,
        "default_art": "constellation",
        "palette": "neon",
        "overlay_opacity": 0.10,
    },
    "nature": {
        "bg_color": (242, 238, 228),
        "title_color": (25, 55, 25),
        "body_color": (60, 80, 45),
        "accent_color": (40, 130, 80),
        "accent2_color": (80, 170, 90),
        "font_title": "Arial",
        "font_body": "Arial",
        "title_size": 40,
        "body_size": 16,
        "default_art": "fractal_tree",
        "palette": "forest",
        "overlay_opacity": 0.08,
    },
    "neon": {
        "bg_color": (4, 4, 8),
        "title_color": (255, 0, 128),
        "body_color": (195, 195, 205),
        "accent_color": (0, 255, 255),
        "accent2_color": (255, 0, 200),
        "font_title": "Arial",
        "font_body": "Arial",
        "title_size": 48,
        "body_size": 18,
        "default_art": "bokeh",
        "palette": "neon",
        "overlay_opacity": 0.15,
    },

    # ─── Chinese Education Themes ─────────────────────────────────────
    "cn_education": {
        "bg_color": (245, 250, 255),
        "title_color": (20, 60, 120),
        "body_color": (50, 50, 60),
        "accent_color": (0, 120, 215),
        "accent2_color": (0, 176, 80),
        "font_title": "Noto Sans CJK SC",
        "font_body": "Noto Sans CJK SC",
        "title_size": 48,
        "body_size": 24,
        "default_art": "wave",
        "palette": "ocean",
        "overlay_opacity": 0.04,
        "task_color": (0, 176, 80),       # Task separator green
        "thought_color": (255, 248, 220),  # Thought bubble cream
        "highlight_color": (255, 255, 0),  # Yellow highlight
    },
    "cn_warm": {
        "bg_color": (255, 250, 240),
        "title_color": (120, 50, 20),
        "body_color": (80, 50, 30),
        "accent_color": (220, 80, 40),
        "accent2_color": (255, 160, 50),
        "font_title": "Noto Serif CJK SC",
        "font_body": "Noto Sans CJK SC",
        "title_size": 48,
        "body_size": 24,
        "default_art": "marble",
        "palette": "warm",
        "overlay_opacity": 0.04,
        "task_color": (220, 80, 40),
        "thought_color": (255, 245, 230),
        "highlight_color": (255, 220, 100),
    },
    "cn_forest": {
        "bg_color": (240, 248, 240),
        "title_color": (20, 70, 30),
        "body_color": (40, 60, 35),
        "accent_color": (30, 130, 60),
        "accent2_color": (80, 180, 80),
        "font_title": "Noto Sans CJK SC",
        "font_body": "Noto Sans CJK SC",
        "title_size": 48,
        "body_size": 24,
        "default_art": "fractal_tree",
        "palette": "forest",
        "overlay_opacity": 0.04,
        "task_color": (30, 130, 60),
        "thought_color": (245, 255, 240),
        "highlight_color": (180, 255, 100),
    },
    "cn_pastel": {
        "bg_color": (252, 248, 255),
        "title_color": (80, 40, 120),
        "body_color": (60, 50, 70),
        "accent_color": (150, 80, 200),
        "accent2_color": (255, 150, 180),
        "font_title": "Noto Sans CJK SC",
        "font_body": "Noto Sans CJK SC",
        "title_size": 48,
        "body_size": 24,
        "default_art": "bokeh",
        "palette": "pastel",
        "overlay_opacity": 0.03,
        "task_color": (150, 80, 200),
        "thought_color": (255, 245, 255),
        "highlight_color": (255, 200, 230),
    },
    "cn_bright": {
        "bg_color": (255, 255, 255),
        "title_color": (0, 80, 180),
        "body_color": (40, 40, 50),
        "accent_color": (0, 150, 220),
        "accent2_color": (255, 90, 70),
        "font_title": "Noto Sans CJK SC",
        "font_body": "Noto Sans CJK SC",
        "title_size": 48,
        "body_size": 24,
        "default_art": "gradient_mesh",
        "palette": "arctic",
        "overlay_opacity": 0.03,
        "task_color": (0, 150, 220),
        "thought_color": (240, 248, 255),
        "highlight_color": (100, 200, 255),
    },
}

# Slide standard dimensions
SLIDE_W_INCH = 13.333
SLIDE_H_INCH = 7.5
SLIDE_W = Inches(SLIDE_W_INCH)
SLIDE_H = Inches(SLIDE_H_INCH)

# Pixel resolution for art generation (at 96 DPI)
ART_W = int(SLIDE_W_INCH * 96)
ART_H = int(SLIDE_H_INCH * 96)

layout = SmartLayout()


# ─── Helpers ───────────────────────────────────────────────────────────────

def _rgb(r, g, b):
    return RGBColor(r, g, b)


def _rgb_tuple(t):
    return RGBColor(*t[:3])


def _save_img(img, fmt="PNG") -> io.BytesIO:
    buf = io.BytesIO()
    img.save(buf, format=fmt)
    buf.seek(0)
    return buf


def _set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _rgb_tuple(color)


def _ensure_pt(val):
    """Ensure value is a Pt object for font size."""
    if isinstance(val, Pt):
        return val
    if isinstance(val, (int, float)):
        if val > 100000:
            val = int(val / 12700)
        return Pt(int(val))
    return Pt(18)


def _get_theme_font(theme_cfg, key="font_title"):
    """Get font name from theme, with fallback."""
    return theme_cfg.get(key, theme_cfg.get("font_body", "Arial"))


def _add_text(slide, left, top, width, height, text, font_size=18,
               color=(255, 255, 255), bold=False, alignment=PP_ALIGN.LEFT,
               font_name="Arial", line_spacing=1.25, italic=False):
    """Add a rich text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    fs = _ensure_pt(font_size)
    fs_val = fs.pt if isinstance(fs, Pt) else 18

    if isinstance(text, list):
        for i, line in enumerate(text):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = fs
            p.font.color.rgb = _rgb_tuple(color)
            p.font.bold = bold
            p.font.name = font_name
            p.font.italic = italic
            p.alignment = alignment
            p.space_after = Pt(fs_val * 0.6)
            if line_spacing and line_spacing != 1.0:
                p.line_spacing = Pt(fs_val * line_spacing)
    else:
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = fs
        p.font.color.rgb = _rgb_tuple(color)
        p.font.bold = bold
        p.font.name = font_name
        p.font.italic = italic
        p.alignment = alignment

    return txBox


def _add_line(slide, left, top, width, color, height_pt=2):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(height_pt))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb_tuple(color)
    shape.line.fill.background()
    return shape


def _add_accent_bar(slide, color, width_emu=Inches(0.07)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width_emu, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb_tuple(color)
    shape.line.fill.background()


def _add_circle(slide, cx, cy, r, color, opacity=1.0):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        cx - r, cy - r, r * 2, r * 2)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb_tuple(color)
    shape.line.fill.background()
    return shape


def _generate_bg_art(slide_config, theme_cfg, art_w=ART_W, art_h=ART_H):
    """Generate background art image."""
    art_style = slide_config.get("bg_style", theme_cfg.get("default_art", "gradient_mesh"))
    palette = slide_config.get("palette", theme_cfg.get("palette", "sunset"))
    seed = slide_config.get("seed", random.randint(1, 999999))
    img = generate_art(art_style, art_w, art_h, palette_name=palette, seed=seed)
    return img, palette


def _apply_bg_art(slide, img, theme_cfg, blur_radius=2, darken=100):
    """Apply art as slide background with optional effects."""
    if blur_radius > 0:
        img = img.filter(ImageFilter.GaussianBlur(radius=blur_radius))
    overlay = Image.new("RGBA", img.size, (0, 0, 0, darken))
    img = img.convert("RGBA")
    img = Image.alpha_composite(img, overlay)
    img = img.convert("RGB")
    buf = _save_img(img)
    slide.shapes.add_picture(buf, 0, 0, SLIDE_W, SLIDE_H)


def _apply_subtle_bg_art(slide, img, theme_cfg, opacity=0.12):
    """Apply art as subtle background texture."""
    img = img.convert("RGBA")
    alpha = img.split()[3] if img.mode == "RGBA" else img.convert("L")
    alpha = alpha.point(lambda p: int(p * opacity))
    img.putalpha(alpha)
    bg = Image.new("RGBA", img.size, theme_cfg["bg_color"] + (255,))
    bg = Image.alpha_composite(bg, img).convert("RGB")
    buf = _save_img(bg)
    slide.shapes.add_picture(buf, 0, 0, SLIDE_W, SLIDE_H)


def _apply_decorative_bg(slide, theme_cfg, seed=42):
    """Apply a subtle decorative background pattern."""
    art_style = theme_cfg.get("default_art", "gradient_mesh")
    palette = theme_cfg.get("palette", "mono")
    opacity = theme_cfg.get("overlay_opacity", 0.05)
    img = generate_art(art_style, ART_W, ART_H, palette_name=palette, seed=seed)
    _apply_subtle_bg_art(slide, img, theme_cfg, opacity=opacity)


def _add_illustration(slide, slide_cfg, theme_cfg, pos=None):
    """Auto-generate and add contextual illustration to a slide."""
    if pos is None:
        pos = (Inches(9.0), Inches(1.5), Inches(3.5), Inches(4.5))
    palette_name = slide_cfg.get("palette", theme_cfg.get("palette", "forest"))
    seed = slide_cfg.get("seed", random.randint(1, 999999))
    try:
        svg_str = auto_illustrate_slide(slide_cfg, palette_name=palette_name, seed=seed)
        import tempfile, os
        fd, png_path = tempfile.mkstemp(suffix=".png")
        os.close(fd)
        _svg_to_png(svg_str, 400, 400, png_path)
        if os.path.exists(png_path) and os.getsize(png_path) > 0:
            buf = io.BytesIO(open(png_path, "rb").read())
            slide.shapes.add_picture(buf, *pos)
            os.unlink(png_path)
    except Exception:
        pass


def _add_decorative_corner(slide, theme_cfg, position="top-right"):
    """Add decorative corner element."""
    accent = theme_cfg.get("accent2_color", theme_cfg.get("accent_color", (0, 120, 215)))
    light_accent = tuple(min(255, c + 80) for c in accent)

    if position == "top-right":
        # Decorative circles in corner
        _add_circle(slide, Inches(12.5), Inches(0.5), Inches(0.3), accent)
        _add_circle(slide, Inches(12.0), Inches(0.8), Inches(0.15), light_accent)
        _add_circle(slide, Inches(12.8), Inches(1.0), Inches(0.1), light_accent)
    elif position == "bottom-left":
        _add_circle(slide, Inches(0.5), Inches(6.8), Inches(0.25), accent)
        _add_circle(slide, Inches(1.0), Inches(7.0), Inches(0.12), light_accent)


# ═══════════════════════════════════════════════════════════════════════════
# ─── ORIGINAL SLIDE BUILDERS (ENHANCED) ────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════

def _build_title(prs, cfg, theme):
    """Full-bleed art title slide with overlay text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme_cfg = THEMES.get(theme, THEMES["dark"])
    art_img, _ = _generate_bg_art(cfg, theme_cfg)
    _apply_bg_art(slide, art_img, theme_cfg, blur_radius=3, darken=120)
    pos = layout.title_position(vertical_center=True)
    heading = cfg.get("heading", "Presentation Title")
    subheading = cfg.get("subheading", "")
    _add_text(slide, *pos["title"], heading,
              font_size=pos["size_title"], color=(255, 255, 255),
              bold=True, alignment=PP_ALIGN.CENTER, font_name=_get_theme_font(theme_cfg))
    _add_line(slide, *pos["divider"], (255, 255, 255), height_pt=2)
    if subheading:
        _add_text(slide, *pos["subtitle"], subheading,
                  font_size=pos["size_subtitle"], color=(210, 210, 230),
                  alignment=PP_ALIGN.CENTER, font_name=_get_theme_font(theme_cfg, "font_body"))


def _build_content(prs, cfg, theme):
    """Content slide with heading, accent bar, and bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme_cfg = THEMES.get(theme, THEMES["dark"])
    art_style = cfg.get("bg_style")
    if art_style:
        art_img, _ = _generate_bg_art(cfg, theme_cfg)
        _apply_subtle_bg_art(slide, art_img, theme_cfg)
    else:
        _set_bg(slide, theme_cfg["bg_color"])
        _apply_decorative_bg(slide, theme_cfg, seed=cfg.get("seed", 42))

    _add_accent_bar(slide, theme_cfg["accent_color"])
    _add_decorative_corner(slide, theme_cfg)
    pos = layout.content_position()
    heading = cfg.get("heading", "Section Title")
    body = cfg.get("body", [])
    _add_text(slide, *pos["heading"], heading,
              font_size=pos["size_heading"], color=theme_cfg["title_color"],
              bold=True, font_name=_get_theme_font(theme_cfg))
    _add_line(slide, *pos["divider"], theme_cfg["accent_color"], height_pt=3)
    if isinstance(body, str):
        body = [body]
    _add_text(slide, *pos["body"], body,
              font_size=pos["size_body"], color=theme_cfg["body_color"],
              font_name=_get_theme_font(theme_cfg, "font_body"))
    if not cfg.get("no_illustration", False):
        ill_pos = pos.get("illustration", (Inches(9.0), Inches(1.5), Inches(3.5), Inches(4.5)))
        _add_illustration(slide, cfg, theme_cfg, ill_pos)


def _build_two_column(prs, cfg, theme):
    """Two-column comparison slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme_cfg = THEMES.get(theme, THEMES["dark"])
    _set_bg(slide, theme_cfg["bg_color"])
    _apply_decorative_bg(slide, theme_cfg, seed=cfg.get("seed", 42))
    pos = layout.two_column_position()
    heading = cfg.get("heading", "Comparison")
    _add_text(slide, *pos["heading"], heading,
              font_size=pos["size_heading"], color=theme_cfg["title_color"],
              bold=True, font_name=_get_theme_font(theme_cfg))
    _add_line(slide, Inches(0.8), Inches(1.25), Inches(2), theme_cfg["accent_color"], 3)
    left = cfg.get("left", {})
    right = cfg.get("right", {})
    _add_text(slide, *pos["left_title"], left.get("title", "Column A"),
              font_size=pos["size_col_title"], color=theme_cfg["accent_color"],
              bold=True, font_name=_get_theme_font(theme_cfg))
    items = left.get("items", [])
    if isinstance(items, str):
        items = [items]
    _add_text(slide, *pos["left_body"], items,
              font_size=pos["size_body"], color=theme_cfg["body_color"],
              font_name=_get_theme_font(theme_cfg, "font_body"))
    div_key = "divider_v" if "divider_v" in pos else "divider"
    div = pos[div_key]
    if len(div) == 4:
        div_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, *div)
        div_shape.fill.solid()
        div_shape.fill.fore_color.rgb = _rgb_tuple(theme_cfg["accent_color"])
        div_shape.line.fill.background()
    else:
        _add_line(slide, *div, theme_cfg["accent_color"], height_pt=1)
    _add_text(slide, *pos["right_title"], right.get("title", "Column B"),
              font_size=pos["size_col_title"], color=theme_cfg["accent_color"],
              bold=True, font_name=_get_theme_font(theme_cfg))
    items = right.get("items", [])
    if isinstance(items, str):
        items = [items]
    _add_text(slide, *pos["right_body"], items,
              font_size=pos["size_body"], color=theme_cfg["body_color"],
              font_name=_get_theme_font(theme_cfg, "font_body"))
    if not cfg.get("no_illustration", False):
        _add_illustration(slide, cfg, theme_cfg,
                         (Inches(10.5), Inches(5.5), Inches(2.2), Inches(1.5)))


def _build_three_column(prs, cfg, theme):
    """Three-column layout slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme_cfg = THEMES.get(theme, THEMES["dark"])
    _set_bg(slide, theme_cfg["bg_color"])
    _apply_decorative_bg(slide, theme_cfg, seed=cfg.get("seed", 42))
    pos = layout.three_column_position()
    heading = cfg.get("heading", "Overview")
    _add_text(slide, *pos["heading"], heading,
              font_size=pos["size_heading"], color=theme_cfg["title_color"],
              bold=True, font_name=_get_theme_font(theme_cfg))
    _add_line(slide, Inches(0.8), Inches(1.25), Inches(2), theme_cfg["accent_color"], 3)
    columns = cfg.get("columns", [{}, {}, {}])
    for i, col_cfg in enumerate(columns[:3]):
        col_pos = pos["columns"][i]
        _add_text(slide, *col_pos["title"], col_cfg.get("title", f"Column {i+1}"),
                  font_size=pos["size_col_title"], color=theme_cfg["accent_color"],
                  bold=True, font_name=_get_theme_font(theme_cfg))
        items = col_cfg.get("items", [])
        if isinstance(items, str):
            items = [items]
        _add_text(slide, *col_pos["body"], items,
                  font_size=pos["size_body"], color=theme_cfg["body_color"],
                  font_name=_get_theme_font(theme_cfg, "font_body"))
    if not cfg.get("no_illustration", False):
        _add_illustration(slide, cfg, theme_cfg,
                         (Inches(10.5), Inches(0.2), Inches(2.2), Inches(1.5)))


def _build_image_full(prs, cfg, theme):
    """Full-bleed art image slide with overlay text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme_cfg = THEMES.get(theme, THEMES["dark"])
    art_img, _ = _generate_bg_art(cfg, theme_cfg, ART_W, ART_H)
    effects = cfg.get("effects", [])
    if effects:
        art_img = apply_effects(art_img, effects)
    w, h = art_img.size
    overlay = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)
    for y in range(int(h * 0.55), h):
        alpha = int(180 * ((y - h * 0.55) / (h * 0.45)))
        draw.line([(0, y), (w, y)], fill=(0, 0, 0, min(180, alpha)))
    art_img = art_img.convert("RGBA")
    art_img = Image.alpha_composite(art_img, overlay).convert("RGB")
    buf = _save_img(art_img)
    slide.shapes.add_picture(buf, 0, 0, SLIDE_W, SLIDE_H)
    pos = layout.image_full_position()
    heading = cfg.get("heading", "")
    caption = cfg.get("caption", "")
    if heading:
        _add_text(slide, *pos["heading"], heading,
                  font_size=pos["size_heading"], color=(255, 255, 255),
                  bold=True, alignment=PP_ALIGN.CENTER, font_name=_get_theme_font(theme_cfg))
    if caption:
        _add_text(slide, *pos["caption"], caption,
                  font_size=pos["size_caption"], color=(200, 200, 210),
                  alignment=PP_ALIGN.CENTER, italic=True)


def _build_chart(prs, cfg, theme):
    """Data chart slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme_cfg = THEMES.get(theme, THEMES["dark"])
    _set_bg(slide, theme_cfg["bg_color"])
    _apply_decorative_bg(slide, theme_cfg, seed=cfg.get("seed", 42))
    pos = layout.chart_position()
    heading = cfg.get("heading", "Data")
    _add_text(slide, *pos["heading"], heading,
              font_size=pos["size_heading"], color=theme_cfg["title_color"],
              bold=True, font_name=_get_theme_font(theme_cfg))
    _add_line(slide, Inches(0.8), Inches(1.25), Inches(2), theme_cfg["accent_color"], 3)
    chart_type_str = cfg.get("chart_type", "bar")
    chart_type_map = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "area": XL_CHART_TYPE.AREA,
        "pie": XL_CHART_TYPE.PIE,
        "donut": XL_CHART_TYPE.DOUGHNUT,
        "scatter": XL_CHART_TYPE.XY_SCATTER,
    }
    data = cfg.get("data", {"labels": ["A", "B", "C"], "values": [30, 50, 20]})
    labels = data.get("labels", [])
    values = data.get("values", [])
    from pptx.chart.data import CategoryChartData
    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series(cfg.get("series_name", "Series 1"), values)
    series_list = data.get("series", [])
    for s in series_list:
        if isinstance(s, dict):
            chart_data.add_series(s.get("name", ""), s.get("values", []))
    chart_shape = slide.shapes.add_chart(
        chart_type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED),
        *pos["chart"], chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = cfg.get("show_legend", True)
    plot = chart.plots[0]
    plot.gap_width = 100
    palette_name = cfg.get("palette", theme_cfg.get("palette", "sunset"))
    palette = get_palette(palette_name)
    for i, series in enumerate(plot.series):
        color = palette[i % len(palette)]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = _rgb_tuple(color)


def _build_end(prs, cfg, theme):
    """Thank you / closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme_cfg = THEMES.get(theme, THEMES["dark"])
    art_img, _ = _generate_bg_art(cfg, theme_cfg)
    _apply_bg_art(slide, art_img, theme_cfg, blur_radius=3, darken=120)
    pos = layout.end_position()
    heading = cfg.get("heading", "Thank You")
    subheading = cfg.get("subheading", "")
    title_key = "title" if "title" in pos else "heading"
    sub_key = "subtitle" if "subtitle" in pos else "subheading"
    _add_text(slide, *pos[title_key], heading,
              font_size=pos["size_title"] if "size_title" in pos else pos["size_heading"],
              color=(255, 255, 255),
              bold=True, alignment=PP_ALIGN.CENTER, font_name=_get_theme_font(theme_cfg))
    _add_line(slide, *pos["divider"], (255, 255, 255), height_pt=2)
    if subheading:
        _add_text(slide, *pos[sub_key], subheading,
                  font_size=pos["size_subtitle"] if "size_subtitle" in pos else pos.get("size_subheading", 22),
                  color=(200, 200, 220),
                  alignment=PP_ALIGN.CENTER)


def _build_agenda(prs, cfg, theme):
    """Agenda / outline slide with numbered items."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme_cfg = THEMES.get(theme, THEMES["dark"])
    _set_bg(slide, theme_cfg["bg_color"])
    _apply_decorative_bg(slide, theme_cfg, seed=cfg.get("seed", 42))
    pos = layout.agenda_position()
    heading = cfg.get("heading", "Agenda")
    _add_text(slide, *pos["heading"], heading,
              font_size=pos["size_heading"], color=theme_cfg["title_color"],
              bold=True, font_name=_get_theme_font(theme_cfg))
    _add_line(slide, Inches(0.8), Inches(1.25), Inches(2), theme_cfg["accent_color"], 3)
    items = cfg.get("items", [])
    palette = get_palette(cfg.get("palette", theme_cfg.get("palette", "sunset")))
    for i, item in enumerate(items):
        y = Inches(1.8) + Inches(i * 0.75)
        number_color = palette[(i + 1) % len(palette)]
        _add_circle(slide, Inches(1.5), y + Inches(0.18), Inches(0.22), number_color)
        _add_text(slide, Inches(1.2), y, Inches(0.6), Inches(0.5),
                  str(i + 1), font_size=pos["size_number"],
                  color=(255, 255, 255), bold=True, alignment=PP_ALIGN.CENTER)
        _add_text(slide, Inches(2.1), y, Inches(9), Inches(0.5),
                  item, font_size=pos["size_item"],
                  color=theme_cfg["body_color"], font_name=_get_theme_font(theme_cfg, "font_body"))
    if not cfg.get("no_illustration", False):
        _add_illustration(slide, cfg, theme_cfg,
                         (Inches(10.2), Inches(5.5), Inches(2.5), Inches(1.5)))


def _build_timeline(prs, cfg, theme):
    """Timeline / process flow slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme_cfg = THEMES.get(theme, THEMES["dark"])
    _set_bg(slide, theme_cfg["bg_color"])
    _apply_decorative_bg(slide, theme_cfg, seed=cfg.get("seed", 42))
    pos = layout.timeline_position()
    heading = cfg.get("heading", "Timeline")
    _add_text(slide, *pos["heading"], heading,
              font_size=pos["size_heading"], color=theme_cfg["title_color"],
              bold=True, font_name=_get_theme_font(theme_cfg))
    _add_line(slide, Inches(0.8), Inches(1.15), Inches(2), theme_cfg["accent_color"], 3)
    items = cfg.get("items", [])
    palette = get_palette(cfg.get("palette", theme_cfg.get("palette", "sunset")))
    n = len(items)
    if n > 0:
        line_y = Inches(4.0)
        start_x = Inches(1.5)
        end_x = Inches(11.8)
        _add_line(slide, start_x, line_y, end_x - start_x, theme_cfg["body_color"], 2)
        for i, item in enumerate(items):
            t = i / max(1, n - 1)
            x = start_x + (end_x - start_x) * t
            dot_color = palette[(i + 1) % len(palette)]
            _add_circle(slide, x, line_y - Inches(0.1), Inches(0.15), dot_color)
            label = item.get("date", item.get("label", f"Step {i+1}"))
            _add_text(slide, x - Inches(0.8), line_y - Inches(0.8), Inches(1.6), Inches(0.5),
                      label, font_size=pos["size_label"],
                      color=theme_cfg["accent_color"], bold=True, alignment=PP_ALIGN.CENTER)
            title = item.get("title", "")
            if title:
                _add_text(slide, x - Inches(0.9), line_y + Inches(0.3), Inches(1.8), Inches(0.8),
                          title, font_size=pos["size_title"],
                          color=theme_cfg["body_color"], alignment=PP_ALIGN.CENTER)
    if not cfg.get("no_illustration", False):
        _add_illustration(slide, cfg, theme_cfg,
                         (Inches(0.3), Inches(5.5), Inches(2.5), Inches(1.5)))


def _build_quote(prs, cfg, theme):
    """Quote / testimonial slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme_cfg = THEMES.get(theme, THEMES["dark"])
    art_style = cfg.get("bg_style")
    if art_style:
        art_img, _ = _generate_bg_art(cfg, theme_cfg)
        _apply_subtle_bg_art(slide, art_img, theme_cfg, opacity=0.08)
    else:
        _set_bg(slide, theme_cfg["bg_color"])
    pos = layout.quote_position()
    _add_text(slide, Inches(0.5), Inches(1.0), Inches(1.5), Inches(1.5),
              "\u201C", font_size=120, color=theme_cfg["accent_color"],
              alignment=PP_ALIGN.LEFT)
    quote_text = cfg.get("text", "Your quote here")
    author = cfg.get("author", "")
    _add_text(slide, *pos["quote"], quote_text,
              font_size=pos["size_quote"], color=theme_cfg["title_color"],
              bold=False, alignment=PP_ALIGN.LEFT, italic=True,
              font_name=_get_theme_font(theme_cfg))
    _add_line(slide, *pos["divider"], theme_cfg["accent_color"], 2)
    if author:
        _add_text(slide, *pos["author"], f"— {author}",
                  font_size=pos["size_author"], color=theme_cfg["body_color"],
                  alignment=PP_ALIGN.CENTER)


def _build_data_table(prs, cfg, theme):
    """Data table slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme_cfg = THEMES.get(theme, THEMES["dark"])
    _set_bg(slide, theme_cfg["bg_color"])
    _apply_decorative_bg(slide, theme_cfg, seed=cfg.get("seed", 42))
    pos = layout.data_table_position()
    heading = cfg.get("heading", "Data Table")
    _add_text(slide, *pos["heading"], heading,
              font_size=pos["size_heading"], color=theme_cfg["title_color"],
              bold=True, font_name=_get_theme_font(theme_cfg))
    _add_line(slide, Inches(0.8), Inches(1.25), Inches(2), theme_cfg["accent_color"], 3)
    headers = cfg.get("headers", [])
    rows = cfg.get("rows", [])
    if headers and rows:
        add_styled_table(slide, *pos["table"], headers, rows,
                         font_name=_get_theme_font(theme_cfg, "font_body"))


def _build_metrics(prs, cfg, theme):
    """Big number / KPI metrics slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme_cfg = THEMES.get(theme, THEMES["dark"])
    _set_bg(slide, theme_cfg["bg_color"])
    _apply_decorative_bg(slide, theme_cfg, seed=cfg.get("seed", 42))
    pos = layout.metrics_position()
    heading = cfg.get("heading", "Key Metrics")
    _add_text(slide, *pos["heading"], heading,
              font_size=pos["size_heading"], color=theme_cfg["title_color"],
              bold=True, font_name=_get_theme_font(theme_cfg))
    _add_line(slide, Inches(0.8), Inches(1.1), Inches(2), theme_cfg["accent_color"], 3)
    metrics = cfg.get("metrics", [])
    n = len(metrics)
    palette = get_palette(cfg.get("palette", theme_cfg.get("palette", "sunset")))
    col_width = Inches(11.3 / max(1, n))
    start_x = Inches(1.0)
    for i, m in enumerate(metrics):
        cx = start_x + col_width * i + col_width / 2
        number_color = palette[(i + 1) % len(palette)]
        _add_text(slide, cx - Inches(1.5), Inches(2.5), Inches(3), Inches(2),
                  m.get("value", "0"),
                  font_size=pos["size_number"], color=number_color,
                  bold=True, alignment=PP_ALIGN.CENTER,
                  font_name=_get_theme_font(theme_cfg))
        _add_text(slide, cx - Inches(1.5), Inches(4.8), Inches(3), Inches(0.8),
                  m.get("label", ""),
                  font_size=pos["size_label"], color=theme_cfg["body_color"],
                  alignment=PP_ALIGN.CENTER, font_name=_get_theme_font(theme_cfg, "font_body"))
        sub = m.get("sublabel", "")
        if sub:
            _add_text(slide, cx - Inches(1.5), Inches(5.5), Inches(3), Inches(0.6),
                      sub, font_size=TYPE_SCALE["sm"], color=theme_cfg["body_color"],
                      alignment=PP_ALIGN.CENTER, italic=True)


def _build_image_grid(prs, cfg, theme):
    """Image grid layout (2x2 or 3x2)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme_cfg = THEMES.get(theme, THEMES["dark"])
    _set_bg(slide, theme_cfg["bg_color"])
    heading = cfg.get("heading", "Gallery")
    _add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
              heading, font_size=TYPE_SCALE["2xl"], color=theme_cfg["title_color"],
              bold=True, font_name=_get_theme_font(theme_cfg))
    _add_line(slide, Inches(0.8), Inches(1.05), Inches(2), theme_cfg["accent_color"], 3)
    art_styles = cfg.get("art_styles", ["gradient_mesh", "flow_field", "voronoi", "bokeh"])
    palette = cfg.get("palette", theme_cfg.get("palette", "sunset"))
    seed = cfg.get("seed", random.randint(1, 999999))
    cols = min(3, len(art_styles))
    rows = (len(art_styles) + cols - 1) // cols
    img_w = int(10.5 / cols * 96)
    img_h = int(5.0 / rows * 96)
    for i, style in enumerate(art_styles[:6]):
        col = i % cols
        row = i // cols
        x = Inches(1.0 + col * (10.5 / cols) + 0.15)
        y = Inches(1.4 + row * (5.0 / rows) + 0.15)
        w = Inches(10.5 / cols - 0.3)
        h = Inches(5.0 / rows - 0.3)
        art = generate_art(style, img_w, img_h, palette_name=palette, seed=seed + i)
        art = art.filter(ImageFilter.GaussianBlur(radius=1))
        buf = _save_img(art)
        slide.shapes.add_picture(buf, x, y, w, h)


# ═══════════════════════════════════════════════════════════════════════════
# ─── CHINESE EDUCATION SLIDE BUILDERS (NEW) ────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════

def _build_cn_title(prs, cfg, theme):
    """Chinese educational title slide — large characters, decorative elements."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme_cfg = THEMES.get(theme, THEMES.get("cn_education", THEMES["dark"]))

    # Full-bleed background art
    art_img, _ = _generate_bg_art(cfg, theme_cfg)
    _apply_bg_art(slide, art_img, theme_cfg, blur_radius=4, darken=80)

    font = _get_theme_font(theme_cfg)

    heading = cfg.get("heading", "课程标题")
    subheading = cfg.get("subheading", "")

    # Main title - large centered
    _add_text(slide, Inches(1.5), Inches(2), Inches(10.3), Inches(3),
              heading,
              font_size=cfg.get("font_size", TYPE_SCALE["cn_massive"]),
              color=(255, 255, 255), bold=True,
              alignment=PP_ALIGN.CENTER, font_name=font)

    # Decorative divider
    _add_line(slide, Inches(5), Inches(5.2), Inches(3.3),
              theme_cfg.get("accent2_color", (255, 255, 255)), height_pt=3)

    if subheading:
        _add_text(slide, Inches(2.5), Inches(5.5), Inches(8.3), Inches(1),
                  subheading,
                  font_size=cfg.get("subtitle_size", TYPE_SCALE["cn_subtitle"]),
                  color=(220, 220, 240),
                  alignment=PP_ALIGN.CENTER, font_name=font)


def _build_cn_task_separator(prs, cfg, theme):
    """Full-page task/section separator — Chinese education style."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme_cfg = THEMES.get(theme, THEMES.get("cn_education", THEMES["dark"]))

    task_color = theme_cfg.get("task_color", theme_cfg.get("accent2_color", (0, 176, 80)))
    font = _get_theme_font(theme_cfg)

    # Decorative background images (top and bottom strips)
    seed = cfg.get("seed", random.randint(1, 999999))
    palette = cfg.get("palette", theme_cfg.get("palette", "ocean"))

    # Top decoration strip
    top_art = generate_art("wave", ART_W, int(ART_H * 0.25), palette_name=palette, seed=seed)
    top_art = top_art.filter(ImageFilter.GaussianBlur(radius=2))
    buf = _save_img(top_art)
    slide.shapes.add_picture(buf, Inches(1.5), 0, Inches(10.3), Inches(1.8))

    # Bottom decoration strip
    bot_art = generate_art("wave", ART_W, int(ART_H * 0.25), palette_name=palette, seed=seed + 1)
    bot_art = bot_art.filter(ImageFilter.GaussianBlur(radius=2))
    buf = _save_img(bot_art)
    slide.shapes.add_picture(buf, Inches(1.5), Inches(5.5), Inches(10.3), Inches(2))

    # Center decorative image
    center_art = generate_art("bokeh", int(ART_W * 0.4), int(ART_H * 0.35),
                               palette_name=palette, seed=seed + 2)
    center_art = center_art.filter(ImageFilter.GaussianBlur(radius=3))
    buf = _save_img(center_art)
    slide.shapes.add_picture(buf, Inches(5), Inches(2.5), Inches(5), Inches(3))

    # Main text
    text = cfg.get("text", cfg.get("heading", "任务"))
    _add_text(slide, Inches(2), Inches(2.2), Inches(9.3), Inches(3),
              text,
              font_size=cfg.get("font_size", TYPE_SCALE["cn_huge"]),
              color=(255, 255, 255), bold=True,
              alignment=PP_ALIGN.CENTER, font_name=font)


def _build_cn_content(prs, cfg, theme):
    """Chinese educational content — with vertical side label and rich layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme_cfg = THEMES.get(theme, THEMES.get("cn_education", THEMES["dark"]))

    _set_bg(slide, theme_cfg["bg_color"])
    _apply_decorative_bg(slide, theme_cfg, seed=cfg.get("seed", 42))

    font = _get_theme_font(theme_cfg)
    accent = theme_cfg.get("accent_color", (0, 120, 215))

    # Decorative side image (right)
    seed = cfg.get("seed", random.randint(1, 999999))
    palette = cfg.get("palette", theme_cfg.get("palette", "ocean"))
    side_art = generate_art("bokeh", int(ART_W * 0.2), int(ART_H * 0.3),
                             palette_name=palette, seed=seed + 5)
    side_art = side_art.filter(ImageFilter.GaussianBlur(radius=2))
    buf = _save_img(side_art)
    slide.shapes.add_picture(buf, Inches(9.5), Inches(4.5), Inches(3.5), Inches(2.8))

    # Vertical side label (left)
    side_label = cfg.get("side_label", "")
    if side_label:
        # Label background
        label_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.2), Inches(1.2), Inches(1.3), Inches(5)
        )
        label_bg.fill.solid()
        label_bg.fill.fore_color.rgb = _rgb_tuple(accent)
        label_bg.line.fill.background()

        add_vertical_text(slide, Inches(0.35), Inches(1.5), Inches(1), Inches(4.5),
                          side_label, font_size=TYPE_SCALE["cn_title"],
                          color=(255, 255, 255), bold=True, font_name=font)

    # Content frame
    content_left = Inches(1.8) if side_label else Inches(0.5)
    content_w = Inches(7.5) if side_label else Inches(11)

    # Framed content area
    frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        content_left, Inches(0.8), content_w, Inches(6.2)
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = _rgb(255, 255, 255)
    frame.line.color.rgb = _rgb(220, 225, 235)
    frame.line.width = Pt(1.5)

    # Heading
    heading = cfg.get("heading", "")
    if heading:
        _add_text(slide, content_left + Inches(0.3), Inches(1), content_w - Inches(0.6), Inches(0.7),
                  heading, font_size=TYPE_SCALE["cn_heading"],
                  color=theme_cfg["title_color"], bold=True, font_name=font)
        _add_line(slide, content_left + Inches(0.3), Inches(1.75),
                  Inches(2.5), accent, height_pt=3)

    # Body
    body = cfg.get("body", [])
    if isinstance(body, str):
        body = [body]
    if body:
        _add_text(slide, content_left + Inches(0.3), Inches(2), content_w - Inches(0.6), Inches(4.5),
                  body, font_size=TYPE_SCALE["cn_body"],
                  color=theme_cfg["body_color"], font_name=font)

    # Illustration
    if not cfg.get("no_illustration", False):
        _add_illustration(slide, cfg, theme_cfg,
                         (content_left + content_w - Inches(3.8), Inches(4),
                          Inches(3.5), Inches(2.8)))

    _add_decorative_corner(slide, theme_cfg)


def _build_cn_scenario(prs, cfg, theme):
    """Chinese education scenario slide — situation + thought bubble + analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme_cfg = THEMES.get(theme, THEMES.get("cn_education", THEMES["dark"]))

    _set_bg(slide, theme_cfg["bg_color"])
    _apply_decorative_bg(slide, theme_cfg, seed=cfg.get("seed", 42))

    font = _get_theme_font(theme_cfg)
    accent = theme_cfg.get("accent_color", (0, 120, 215))
    thought_bg = theme_cfg.get("thought_color", (255, 248, 220))

    # Scenario label (top-left badge)
    scenario_label = cfg.get("scenario_label", "情境")
    label_bg_color = cfg.get("label_color", accent)

    label_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.2), Inches(0), Inches(2.5), Inches(1.1)
    )
    label_shape.fill.solid()
    label_shape.fill.fore_color.rgb = _rgb_tuple(label_bg_color)
    label_shape.line.fill.background()

    _add_text(slide, Inches(0.5), Inches(0.15), Inches(2), Inches(0.8),
              scenario_label,
              font_size=TYPE_SCALE["cn_heading"],
              color=(255, 255, 255), bold=True, font_name=font)

    # Content box (right side)
    content_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4.8), Inches(0.8), Inches(8.3), Inches(6.2)
    )
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = _rgb(255, 255, 255)
    content_box.line.color.rgb = _rgb(220, 225, 235)
    content_box.line.width = Pt(1.5)

    # Scenario text
    scenario_text = cfg.get("text", "")
    if scenario_text:
        _add_text(slide, Inches(5.2), Inches(1.2), Inches(7.5), Inches(1.8),
                  scenario_text,
                  font_size=TYPE_SCALE["cn_body"],
                  color=theme_cfg["title_color"], bold=True, font_name=font)

    # Connector lines (decorative)
    for x_offset, length in [(0.5, 3.5), (1.0, 2.8), (1.5, 2.0)]:
        _add_line(slide, Inches(5 + x_offset), Inches(3.2),
                  Inches(length), accent, height_pt=1)

    # Image area
    image_cfg = cfg.get("image", {})
    if image_cfg.get("path"):
        import os
        if os.path.exists(image_cfg["path"]):
            with open(image_cfg["path"], 'rb') as f:
                buf = io.BytesIO(f.read())
            slide.shapes.add_picture(buf,
                Inches(image_cfg.get("left", 6)), Inches(image_cfg.get("top", 3.5)),
                Inches(image_cfg.get("width", 5)), Inches(image_cfg.get("height", 3)))

    # Thought bubble (bottom-left)
    thought_text = cfg.get("thought", "")
    if thought_text:
        add_thought_bubble(slide,
            Inches(0.2), Inches(4.5), Inches(4.3), Inches(2.5),
            text=thought_text,
            font_size=TYPE_SCALE["cn_label"],
            text_color=theme_cfg["title_color"],
            fill_color=thought_bg,
            border_color=accent,
            font_name=font)

    # Decorative corner image
    seed = cfg.get("seed", random.randint(1, 999999))
    palette = cfg.get("palette", theme_cfg.get("palette", "ocean"))
    corner_art = generate_art("bokeh", int(ART_W * 0.15), int(ART_H * 0.2),
                               palette_name=palette, seed=seed + 3)
    corner_art = corner_art.filter(ImageFilter.GaussianBlur(radius=2))
    buf = _save_img(corner_art)
    slide.shapes.add_picture(buf, Inches(9.5), Inches(5), Inches(3.5), Inches(2.2))


def _build_cn_mind_map(prs, cfg, theme):
    """Mind map layout for knowledge organization."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme_cfg = THEMES.get(theme, THEMES.get("cn_education", THEMES["dark"]))

    _set_bg(slide, theme_cfg["bg_color"])
    _apply_decorative_bg(slide, theme_cfg, seed=cfg.get("seed", 42))

    font = _get_theme_font(theme_cfg)

    # Heading
    heading = cfg.get("heading", "思维导图")
    _add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
              heading, font_size=TYPE_SCALE["cn_heading"],
              color=theme_cfg["title_color"], bold=True, font_name=font)
    _add_line(slide, Inches(0.8), Inches(1.05), Inches(2),
              theme_cfg["accent_color"], height_pt=3)

    # Mind map
    root_text = cfg.get("root_text", "中心主题")
    children = cfg.get("children", [])
    root_pos = cfg.get("root_pos", (Inches(5.5), Inches(4)))

    root_color = theme_cfg.get("accent_color", (0, 176, 80))
    child_colors = [
        theme_cfg.get("accent2_color", (0, 120, 215)),
        (255, 90, 70), (255, 160, 0), (128, 0, 200),
        (0, 180, 150), (200, 50, 100)
    ]

    add_mind_map(slide, root_text, children,
                 root_pos=root_pos,
                 root_color=root_color,
                 child_colors=child_colors,
                 font_name=font)

    # Decorative
    seed = cfg.get("seed", random.randint(1, 999999))
    palette = cfg.get("palette", theme_cfg.get("palette", "ocean"))
    corner_art = generate_art("geometric", int(ART_W * 0.12), int(ART_H * 0.15),
                               palette_name=palette, seed=seed + 7)
    corner_art = corner_art.filter(ImageFilter.GaussianBlur(radius=3))
    buf = _save_img(corner_art)
    slide.shapes.add_picture(buf, Inches(9.8), Inches(5.5), Inches(3.2), Inches(1.8))


def _build_cn_evaluation(prs, cfg, theme):
    """Evaluation / rubric table — Chinese education style."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme_cfg = THEMES.get(theme, THEMES.get("cn_education", THEMES["dark"]))

    _set_bg(slide, theme_cfg["bg_color"])
    _apply_decorative_bg(slide, theme_cfg, seed=cfg.get("seed", 42))

    font = _get_theme_font(theme_cfg)

    # Heading
    heading = cfg.get("heading", "评价表")
    _add_text(slide, Inches(1.5), Inches(0.4), Inches(10), Inches(0.8),
              heading, font_size=TYPE_SCALE["cn_heading"],
              color=theme_cfg["title_color"], bold=True, font_name=font,
              alignment=PP_ALIGN.CENTER)

    # Evaluation table
    headers = cfg.get("headers", [])
    rows = cfg.get("rows", [])
    if headers and rows:
        header_bg = cfg.get("header_color", theme_cfg.get("accent_color", (0, 120, 215)))
        add_styled_table(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5),
                         headers, rows,
                         header_bg=header_bg,
                         font_size=TYPE_SCALE["cn_caption"],
                         font_name=font)

    # Decorative image (bottom-left)
    seed = cfg.get("seed", random.randint(1, 999999))
    palette = cfg.get("palette", theme_cfg.get("palette", "ocean"))
    deco_art = generate_art("bokeh", int(ART_W * 0.12), int(ART_H * 0.2),
                             palette_name=palette, seed=seed + 10)
    deco_art = deco_art.filter(ImageFilter.GaussianBlur(radius=3))
    buf = _save_img(deco_art)
    slide.shapes.add_picture(buf, Inches(0.1), Inches(4.5), Inches(1), Inches(2.5))


def _build_cn_flowchart(prs, cfg, theme):
    """Process flow / writing guide layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme_cfg = THEMES.get(theme, THEMES.get("cn_education", THEMES["dark"]))

    _set_bg(slide, theme_cfg["bg_color"])
    _apply_decorative_bg(slide, theme_cfg, seed=cfg.get("seed", 42))

    font = _get_theme_font(theme_cfg)
    accent = theme_cfg.get("accent_color", (0, 120, 215))

    # Heading
    heading = cfg.get("heading", "流程图")
    _add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
              heading, font_size=TYPE_SCALE["cn_heading"],
              color=theme_cfg["title_color"], bold=True, font_name=font)
    _add_line(slide, Inches(0.8), Inches(1.05), Inches(2), accent, height_pt=3)

    # Steps
    steps = cfg.get("steps", [])
    palette = get_palette(cfg.get("palette", theme_cfg.get("palette", "ocean")))
    n = len(steps)
    if n > 0:
        # Calculate layout: up to 4 per row
        per_row = min(4, n)
        rows_count = (n + per_row - 1) // per_row
        box_w = Inches(2.8)
        box_h = Inches(1.8)
        gap_x = Inches(0.3)
        gap_y = Inches(0.5)
        start_x = Inches(0.5 + (12.3 - per_row * 2.8 - (per_row - 1) * 0.3) / 2)

        for i, step in enumerate(steps):
            row = i // per_row
            col = i % per_row
            x = start_x + col * (box_w + gap_x)
            y = Inches(1.5) + row * (box_h + gap_y)

            color = palette[(i + 1) % len(palette)]
            add_step_box(slide, x, y, box_w, box_h,
                         number=i + 1,
                         title=step.get("title", f"步骤 {i+1}"),
                         description=step.get("description", ""),
                         accent_color=color,
                         bg_color=tuple(min(255, c + 200) for c in color),
                         font_name=font)

            # Arrow between steps (not after last in row or last overall)
            if col < per_row - 1 and i < n - 1:
                arrow_x = x + box_w + Inches(0.05)
                arrow_y = y + box_h / 2 - Inches(0.1)
                add_right_arrow_shape(slide, arrow_x, arrow_y,
                                     Inches(0.2), Inches(0.2), color=color)

    # Decorative
    seed = cfg.get("seed", random.randint(1, 999999))
    palette_name = cfg.get("palette", theme_cfg.get("palette", "ocean"))
    deco_art = generate_art("geometric", int(ART_W * 0.12), int(ART_H * 0.15),
                             palette_name=palette_name, seed=seed + 12)
    deco_art = deco_art.filter(ImageFilter.GaussianBlur(radius=3))
    buf = _save_img(deco_art)
    slide.shapes.add_picture(buf, Inches(9.8), Inches(5), Inches(3.2), Inches(2.2))


def _build_cn_reading(prs, cfg, theme):
    """Reading passage with analysis sidebar."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme_cfg = THEMES.get(theme, THEMES.get("cn_education", THEMES["dark"]))

    _set_bg(slide, theme_cfg["bg_color"])
    _apply_decorative_bg(slide, theme_cfg, seed=cfg.get("seed", 42))

    font = _get_theme_font(theme_cfg)
    accent = theme_cfg.get("accent_color", (0, 120, 215))

    # Title bar
    title = cfg.get("title", "")
    if title:
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(0.3), Inches(8), Inches(0.7)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = _rgb_tuple(accent)
        title_bar.line.fill.background()

        _add_text(slide, Inches(0.7), Inches(0.35), Inches(7.5), Inches(0.6),
                  title, font_size=TYPE_SCALE["cn_heading"],
                  color=(255, 255, 255), bold=True, font_name=font)

    # Reading passage (left)
    passage = cfg.get("passage", "")
    if passage:
        # White card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.2), Inches(8.3), Inches(5.8)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(255, 255, 255)
        card.line.color.rgb = _rgb(220, 225, 235)
        card.line.width = Pt(1)

        _add_text(slide, Inches(0.8), Inches(1.5), Inches(7.7), Inches(5.2),
                  passage, font_size=TYPE_SCALE["cn_body"],
                  color=theme_cfg["body_color"], font_name=font)

    # Analysis sidebar (right)
    sidebar_items = cfg.get("sidebar", [])
    if sidebar_items:
        sidebar_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(9), Inches(0.8), Inches(4), Inches(6.2)
        )
        sidebar_bg.fill.solid()
        sidebar_bg.fill.fore_color.rgb = _rgb(245, 248, 255)
        sidebar_bg.line.color.rgb = _rgb(200, 210, 230)
        sidebar_bg.line.width = Pt(1)

        _add_text(slide, Inches(9.2), Inches(1), Inches(3.6), Inches(0.5),
                  "分析要点", font_size=TYPE_SCALE["cn_subtitle"],
                  color=accent, bold=True, font_name=font)

        for i, item in enumerate(sidebar_items):
            y = Inches(1.6 + i * 0.6)
            # Bullet
            _add_circle(slide, Inches(9.5), y + Inches(0.12), Inches(0.08),
                       accent)
            _add_text(slide, Inches(9.7), y, Inches(3.1), Inches(0.5),
                      item, font_size=TYPE_SCALE["cn_label"],
                      color=theme_cfg["body_color"], font_name=font)

    # Decorative
    seed = cfg.get("seed", random.randint(1, 999999))
    palette = cfg.get("palette", theme_cfg.get("palette", "ocean"))
    deco_art = generate_art("bokeh", int(ART_W * 0.15), int(ART_H * 0.2),
                             palette_name=palette, seed=seed + 15)
    deco_art = deco_art.filter(ImageFilter.GaussianBlur(radius=3))
    buf = _save_img(deco_art)
    slide.shapes.add_picture(buf, Inches(9.2), Inches(5), Inches(3.5), Inches(2.2))


def _build_cn_brackets(prs, cfg, theme):
    """Bracket analysis slide — curly brackets with labeled groups."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme_cfg = THEMES.get(theme, THEMES.get("cn_education", THEMES["dark"]))

    _set_bg(slide, theme_cfg["bg_color"])
    _apply_decorative_bg(slide, theme_cfg, seed=cfg.get("seed", 42))

    font = _get_theme_font(theme_cfg)

    # Heading
    heading = cfg.get("heading", "")
    if heading:
        _add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                  heading, font_size=TYPE_SCALE["cn_heading"],
                  color=theme_cfg["title_color"], bold=True, font_name=font)
        _add_line(slide, Inches(0.8), Inches(1.05), Inches(2),
                  theme_cfg["accent_color"], height_pt=3)

    # Main content image/area (center)
    center_content = cfg.get("center_image")
    if center_content and center_content.get("path"):
        import os
        if os.path.exists(center_content["path"]):
            with open(center_content["path"], 'rb') as f:
                buf = io.BytesIO(f.read())
            slide.shapes.add_picture(buf,
                Inches(center_content.get("left", 4.5)),
                Inches(center_content.get("top", 1.5)),
                Inches(center_content.get("width", 5)),
                Inches(center_content.get("height", 4)))

    # Bracket groups
    groups = cfg.get("groups", [])
    palette = get_palette(cfg.get("palette", theme_cfg.get("palette", "ocean")))
    for i, group in enumerate(groups):
        color = palette[(i + 1) % len(palette)]
        pos = group.get("position", {})
        add_bracket_group(slide,
            Inches(pos.get("left", 0.3)),
            Inches(pos.get("top", 1.5 + i * 2)),
            Inches(pos.get("width", 4)),
            Inches(pos.get("height", 1.8)),
            label_text=group.get("label", ""),
            items=group.get("items", []),
            label_color=color,
            font_name=font)

    # Decorative
    seed = cfg.get("seed", random.randint(1, 999999))
    palette_name = cfg.get("palette", theme_cfg.get("palette", "ocean"))
    deco_art = generate_art("geometric", int(ART_W * 0.12), int(ART_H * 0.15),
                             palette_name=palette_name, seed=seed + 18)
    deco_art = deco_art.filter(ImageFilter.GaussianBlur(radius=3))
    buf = _save_img(deco_art)
    slide.shapes.add_picture(buf, Inches(9.5), Inches(5.2), Inches(3.5), Inches(2))


# ─── Slide Type Registry ──────────────────────────────────────────────────

SLIDE_BUILDERS = {
    # Original types
    "title": _build_title,
    "content": _build_content,
    "two_column": _build_two_column,
    "three_column": _build_three_column,
    "image_full": _build_image_full,
    "chart": _build_chart,
    "end": _build_end,
    "agenda": _build_agenda,
    "timeline": _build_timeline,
    "quote": _build_quote,
    "data_table": _build_data_table,
    "metrics": _build_metrics,
    "image_grid": _build_image_grid,
    # Chinese education types
    "cn_title": _build_cn_title,
    "cn_task_separator": _build_cn_task_separator,
    "cn_content": _build_cn_content,
    "cn_scenario": _build_cn_scenario,
    "cn_mind_map": _build_cn_mind_map,
    "cn_evaluation": _build_cn_evaluation,
    "cn_flowchart": _build_cn_flowchart,
    "cn_reading": _build_cn_reading,
    "cn_brackets": _build_cn_brackets,
}


# ─── Main Builder ─────────────────────────────────────────────────────────

def build_ppt(config: Dict, output_path: str) -> str:
    """Build a complete PPTX from config dictionary."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    theme = config.get("theme", "dark")
    slides_config = config.get("slides", [])
    title = config.get("title", "Presentation")

    # Set presentation title
    prs.core_properties.title = title

    for slide_cfg in slides_config:
        slide_type = slide_cfg.get("type", "content")
        builder = SLIDE_BUILDERS.get(slide_type, _build_content)

        # Auto-seed if not provided
        if "seed" not in slide_cfg:
            slide_cfg["seed"] = random.randint(1, 999999)

        builder(prs, slide_cfg, theme)

    prs.save(output_path)
    return output_path
