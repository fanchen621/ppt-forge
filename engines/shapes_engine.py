"""
PPT Forge v2 — Decorative Shapes Engine
Rich visual elements: thought bubbles, cloud frames, brackets,
connector lines, mind map nodes, decorative frames, and more.
"""

import math
import random
from typing import List, Tuple, Dict, Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


def _rgb(color: Tuple[int, int, int]) -> RGBColor:
    return RGBColor(*color[:3])


def _set_shape_fill(shape, color: Tuple[int, int, int], alpha: float = 1.0):
    """Set shape fill color with optional transparency."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    if alpha < 1.0:
        shape.fill.fore_color.brightness = 0  # pptx limitation


def _set_no_line(shape):
    """Remove shape border."""
    shape.line.fill.background()


# ─── Thought Bubble ────────────────────────────────────────────────────────

def add_thought_bubble(slide, left, top, width, height,
                        text="", font_size=18, text_color=(80, 80, 80),
                        fill_color=(255, 255, 255), border_color=None,
                        font_name="Arial"):
    """Add a thought bubble (cloud) shape with text."""
    # Main cloud body
    cloud = slide.shapes.add_shape(
        MSO_SHAPE.CLOUD, left, top, width, height
    )
    _set_shape_fill(cloud, fill_color)
    if border_color:
        cloud.line.color.rgb = _rgb(border_color)
        cloud.line.width = Pt(1.5)
    else:
        _set_no_line(cloud)

    # Small circles below (thought trail)
    trail_y = top + height
    trail_x = left + width * 0.3
    for i, (offset_x, r) in enumerate([(0, Inches(0.12)), (Inches(0.2), Inches(0.08))]):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            trail_x + offset_x, trail_y + Inches(i * 0.12),
            r * 2, r * 2
        )
        _set_shape_fill(dot, fill_color)
        if border_color:
            dot.line.color.rgb = _rgb(border_color)
            dot.line.width = Pt(1)
        else:
            _set_no_line(dot)

    # Text inside cloud
    if text:
        txbox = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(0.15),
            width - Inches(0.4), height - Inches(0.3)
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = _rgb(text_color)
        p.font.name = font_name
        p.font.bold = True

    return cloud


# ─── Decorative Frame (Rounded Rectangle with Header) ──────────────────────

def add_framed_content(slide, left, top, width, height,
                        title="", body_text="",
                        title_color=(255, 255, 255),
                        header_bg=(0, 176, 80),
                        body_bg=(255, 255, 255),
                        border_color=(200, 200, 200),
                        title_size=24, body_size=16,
                        font_name="Arial"):
    """Add a rounded rectangle frame with colored header and content."""
    # Outer frame
    frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    _set_shape_fill(frame, body_bg)
    frame.line.color.rgb = _rgb(border_color)
    frame.line.width = Pt(1.5)

    # Header bar
    header_h = Inches(0.55)
    header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, header_h
    )
    _set_shape_fill(header, header_bg)
    _set_no_line(header)

    # Title text
    if title:
        txbox = slide.shapes.add_textbox(
            left + Inches(0.15), top + Inches(0.05),
            width - Inches(0.3), header_h
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(title_size)
        p.font.color.rgb = _rgb(title_color)
        p.font.bold = True
        p.font.name = font_name

    # Body text
    if body_text:
        txbox = slide.shapes.add_textbox(
            left + Inches(0.2), top + header_h + Inches(0.1),
            width - Inches(0.4), height - header_h - Inches(0.2)
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = body_text
        p.font.size = Pt(body_size)
        p.font.color.rgb = _rgb((60, 60, 60))
        p.font.name = font_name

    return frame


# ─── Bracket (Left/Right Curly Brace with Label) ───────────────────────────

def add_bracket_group(slide, left, top, width, height,
                       label_text="", items=None,
                       label_color=(255, 0, 0), label_size=28,
                       item_color=(80, 80, 80), item_size=20,
                       bracket_color=(150, 150, 150),
                       font_name="Arial"):
    """Add a curly bracket with label and item list."""
    if items is None:
        items = []

    # Left bracket shape
    bracket = slide.shapes.add_shape(
        MSO_SHAPE.LEFT_BRACE,
        left, top, Inches(0.3), height
    )
    _set_shape_fill(bracket, bracket_color)
    _set_no_line(bracket)

    # Label next to bracket
    txbox = slide.shapes.add_textbox(
        left + Inches(0.4), top,
        width - Inches(0.5), Inches(0.5)
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label_text
    p.font.size = Pt(label_size)
    p.font.color.rgb = _rgb(label_color)
    p.font.bold = True
    p.font.name = font_name

    # Items below label
    for i, item in enumerate(items):
        item_y = top + Inches(0.6 + i * 0.45)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left + Inches(0.5), item_y + Inches(0.08),
            Inches(0.1), Inches(0.1)
        )
        _set_shape_fill(dot, label_color)
        _set_no_line(dot)

        txbox = slide.shapes.add_textbox(
            left + Inches(0.7), item_y,
            width - Inches(0.8), Inches(0.4)
        )
        tf = txbox.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(item_size)
        p.font.color.rgb = _rgb(item_color)
        p.font.name = font_name


# ─── Mind Map Node ─────────────────────────────────────────────────────────

def add_mind_map_node(slide, cx, cy, radius, text,
                       fill_color=(0, 176, 80),
                       text_color=(255, 255, 255),
                       font_size=16, bold=True,
                       font_name="Arial"):
    """Add an oval mind map node."""
    node = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        cx - radius, cy - radius,
        radius * 2, radius * 2
    )
    _set_shape_fill(node, fill_color)
    _set_no_line(node)

    # Text
    txbox = slide.shapes.add_textbox(
        cx - radius + Inches(0.05), cy - Inches(0.15),
        radius * 2 - Inches(0.1), Inches(0.35)
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = _rgb(text_color)
    p.font.bold = bold
    p.alignment = PP_ALIGN.CENTER
    p.font.name = font_name

    return node


def add_mind_map_connector(slide, x1, y1, x2, y2,
                            color=(150, 150, 150), width_pt=2):
    """Add a line connector between mind map nodes."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        x1, y1, x2, y2
    )
    connector.line.color.rgb = _rgb(color)
    connector.line.width = Pt(width_pt)
    return connector


def add_mind_map(slide, root_text, children, root_pos=None,
                  root_color=(0, 176, 80), child_colors=None,
                  font_name="Arial"):
    """
    Add a complete mind map with root and children.
    children: list of dicts with 'text', 'sub' (optional list of sub-items)
    """
    if root_pos is None:
        root_pos = (Inches(4), Inches(3.5))
    if child_colors is None:
        child_colors = [(0, 120, 215), (255, 90, 70), (255, 160, 0),
                        (128, 0, 200), (0, 180, 150), (200, 50, 100)]

    rx, ry = root_pos
    root_radius = Inches(0.6)

    # Root node
    add_mind_map_node(slide, rx, ry, root_radius, root_text,
                      fill_color=root_color, font_size=20,
                      font_name=font_name)

    n = len(children)
    for i, child in enumerate(children):
        angle = -90 + (180 / max(1, n - 1)) * i if n > 1 else 0
        angle_rad = math.radians(angle)
        dist = Inches(2.5)
        cx = rx + dist * math.cos(angle_rad)
        cy = ry + dist * math.sin(angle_rad)

        color = child_colors[i % len(child_colors)]

        # Child node
        child_radius = Inches(0.45)
        add_mind_map_node(slide, cx, cy, child_radius,
                          child.get("text", ""),
                          fill_color=color, font_size=14,
                          font_name=font_name)

        # Connector line
        add_mind_map_connector(slide, rx, ry, cx, cy, color=color)

        # Sub-items
        subs = child.get("sub", [])
        for j, sub_text in enumerate(subs):
            sub_angle = angle + (j - len(subs) / 2 + 0.5) * 25
            sub_rad = math.radians(sub_angle)
            sub_dist = Inches(1.3)
            sx = cx + sub_dist * math.cos(sub_rad)
            sy = cy + sub_dist * math.sin(sub_rad)

            sub_radius = Inches(0.3)
            lighter = tuple(min(255, c + 60) for c in color)
            add_mind_map_node(slide, sx, sy, sub_radius, sub_text,
                              fill_color=lighter, text_color=(60, 60, 60),
                              font_size=11, bold=False,
                              font_name=font_name)
            add_mind_map_connector(slide, cx, cy, sx, sy,
                                   color=lighter, width_pt=1.5)


# ─── Section Divider (Task Separator) ──────────────────────────────────────

def add_task_separator(slide, text, slide_width=Inches(13.333),
                        slide_height=Inches(7.5),
                        bg_color=(0, 176, 80),
                        text_color=(255, 255, 255),
                        font_size=60, font_name="Arial"):
    """Add a full-page task/section separator with large centered text."""
    # Background overlay
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, slide_width, slide_height
    )
    _set_shape_fill(bg, bg_color)
    _set_no_line(bg)

    # Text
    txbox = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        slide_width - Inches(2), Inches(2.5)
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = _rgb(text_color)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.name = font_name


# ─── Vertical Text Box ─────────────────────────────────────────────────────

def add_vertical_text(slide, left, top, width, height,
                       text, font_size=36,
                       color=(80, 80, 80), bold=True,
                       font_name="Arial"):
    """Add a vertical text box (for Chinese title strips)."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True

    for i, char in enumerate(text):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = char
        p.font.size = Pt(font_size)
        p.font.color.rgb = _rgb(color)
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(2)

    return txbox


# ─── Decorative Number Badge ───────────────────────────────────────────────

def add_number_badge(slide, cx, cy, number, radius=Inches(0.25),
                      fill_color=(0, 176, 80), text_color=(255, 255, 255),
                      font_size=16, font_name="Arial"):
    """Add a circular number badge."""
    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        cx - radius, cy - radius,
        radius * 2, radius * 2
    )
    _set_shape_fill(badge, fill_color)
    _set_no_line(badge)

    txbox = slide.shapes.add_textbox(
        cx - radius, cy - Inches(0.12),
        radius * 2, Inches(0.3)
    )
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(font_size)
    p.font.color.rgb = _rgb(text_color)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.name = font_name

    return badge


# ─── Arrow Connector ───────────────────────────────────────────────────────

def add_arrow(slide, x1, y1, x2, y2,
               color=(100, 100, 100), width_pt=2):
    """Add a line arrow between two points."""
    connector = slide.shapes.add_connector(
        1,  # straight
        x1, y1, x2, y2
    )
    connector.line.color.rgb = _rgb(color)
    connector.line.width = Pt(width_pt)
    # End arrow
    connector.end_x = x2
    connector.end_y = y2
    return connector


def add_right_arrow_shape(slide, left, top, width, height,
                           color=(0, 176, 80)):
    """Add a right-pointing arrow shape."""
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        left, top, width, height
    )
    _set_shape_fill(arrow, color)
    _set_no_line(arrow)
    return arrow


# ─── Decorative Card (Rounded Rectangle with Shadow Effect) ────────────────

def add_card(slide, left, top, width, height,
              fill_color=(255, 255, 255),
              border_color=(220, 220, 220),
              corner_radius=Inches(0.1)):
    """Add a card-style rounded rectangle."""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    _set_shape_fill(card, fill_color)
    card.line.color.rgb = _rgb(border_color)
    card.line.width = Pt(1)
    return card


# ─── Image Collage Frame ───────────────────────────────────────────────────

def add_image_with_frame(slide, img_path, left, top, width, height,
                          border_color=(255, 255, 255), border_width=Pt(4)):
    """Add an image with decorative border frame."""
    import io
    # Frame background
    frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left - border_width, top - border_width,
        width + border_width * 2, height + border_width * 2
    )
    _set_shape_fill(frame, border_color)
    _set_no_line(frame)

    # Image
    with open(img_path, 'rb') as f:
        buf = io.BytesIO(f.read())
    slide.shapes.add_picture(buf, left, top, width, height)


# ─── Flowchart Step Box ────────────────────────────────────────────────────

def add_step_box(slide, left, top, width, height,
                  number, title, description="",
                  accent_color=(0, 120, 215),
                  bg_color=(240, 245, 255),
                  font_name="Arial"):
    """Add a numbered step box for process flows."""
    # Background
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    _set_shape_fill(box, bg_color)
    box.line.color.rgb = _rgb(accent_color)
    box.line.width = Pt(1.5)

    # Number badge
    add_number_badge(slide,
                     left + Inches(0.35), top + Inches(0.35),
                     number, radius=Inches(0.22),
                     fill_color=accent_color, font_size=14,
                     font_name=font_name)

    # Title
    txbox = slide.shapes.add_textbox(
        left + Inches(0.65), top + Inches(0.1),
        width - Inches(0.75), Inches(0.4)
    )
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.color.rgb = _rgb(accent_color)
    p.font.bold = True
    p.font.name = font_name

    # Description
    if description:
        txbox = slide.shapes.add_textbox(
            left + Inches(0.65), top + Inches(0.45),
            width - Inches(0.75), height - Inches(0.55)
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = description
        p.font.size = Pt(12)
        p.font.color.rgb = _rgb((80, 80, 80))
        p.font.name = font_name


# ─── Evaluation Table Row ──────────────────────────────────────────────────

def add_styled_table(slide, left, top, width, height,
                      headers, rows,
                      header_bg=(0, 120, 215),
                      header_text=(255, 255, 255),
                      odd_bg=(245, 248, 255),
                      even_bg=(255, 255, 255),
                      border_color=(200, 210, 230),
                      font_size=12, font_name="Arial"):
    """Add a beautifully styled table with alternating rows."""
    n_rows = len(rows) + 1
    n_cols = len(headers)

    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = tbl_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(font_size + 2)
        p.font.bold = True
        p.font.color.rgb = _rgb(header_text)
        p.font.name = font_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(header_bg)

    # Data rows
    for i, row in enumerate(rows):
        bg = odd_bg if i % 2 == 0 else even_bg
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(font_size)
            p.font.color.rgb = _rgb((60, 60, 60))
            p.font.name = font_name
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(bg)

    return tbl_shape
